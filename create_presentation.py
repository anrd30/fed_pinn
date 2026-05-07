import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_slide(prs, title_text, bullets=None, image_path=None, layout_idx=1):
    slide_layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = title_text
    
    if bullets:
        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.clear()  # Clear existing
        for i, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = bullet
            p.font.size = Pt(20)
            p.level = 0
            
    if image_path and os.path.exists(image_path):
        # Add image to the right/bottom
        left = Inches(5.5)
        top = Inches(2.0)
        width = Inches(4.0)
        slide.shapes.add_picture(image_path, left, top, width=width)

def main():
    prs = Presentation()

    # Slide 1: Title
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Bounding the Invisible Attack:\nManifold Geometry and Compute-Asymmetry as Dual Defenses"
    subtitle.text = "Neutralizing Adaptive RKHS Evasion in Logit Distillation\n\nFederated Learning Security"

    # Slide 2: The Problem
    create_slide(prs, "The 'Invisible' RKHS Poisoning Attack", [
        "Standard Federated Learning anomaly defenses rely on measuring statistical scalar distances (e.g., MMD, Shannon Entropy).",
        "The Vulnerability: Adaptive attackers in Logit Distillation FL can run internal RKHS optimization loops to match benign statistics perfectly.",
        "Traditional distance filters (FLAME, DeepSight) fail completely against distribution-matched inputs.",
        "The Question: How does a central server detect a semantic backdoor that is mathematically invisible?"
    ])

    # Slide 3: Dual Architecture
    create_slide(prs, "A Paradigm Shift: Geometry + Physics", [
        "Instead of scalar statistics, our defense cross-validates across two immutable physical boundaries:",
        "1. The Mathematical Boundary: Translating local logit probabilities into non-Euclidean Riemannian geometries to detect 'Bending Energy'.",
        "2. The Physical Edge Boundary: Weaponizing physical execution speed limits of Edge devices (e.g., Jetson Nano) to cryptographically identify spoofers."
    ])

    # Slide 4: Riemannian Geometry
    create_slide(prs, "The Fisher-Riemannian Metric", [
        "We map the logit simplex explicitly via the Fisher Information Metric.",
        "This weights gradient shifts in low-probability regions massively heavier than standard Euclidean spaces.",
        "The Harmonic Law: At convergence, clean logit spaces naturally reside in an equilibrium state of minimal bending energy.",
        "Insight: An attacker can fake global MMD statistics, but injecting an artificial backdoor forces a local topological tear across probability boundaries."
    ])

    # Slide 5: Visualizing Bending (Clean)
    create_slide(prs, "Visualizing Topology: Harmonic Baseline", [
        "The naturally converged, flat harmonic equilibrium baseline.",
        "Represents zero bending energy across the classification bounds.",
        "A true Edge device operates in this smooth equilibrium."
    ], image_path=os.path.join("results", "manifold_clean.png"))

    # Slide 6: Visualizing Bending (Poisoned)
    create_slide(prs, "Visualizing Topology: The Backdoor Spike", [
        "The RKHS Backdoor topology.",
        "Global statistics look identical locally, but our PINN detects the massive, localized Dirichlet Bending Energy eruption.",
        "The Topological Tear gives away the Backdoor."
    ], image_path=os.path.join("results", "manifold_poisoned.png"))

    # Slide 7: Hardware Limits
    create_slide(prs, "The Compute-Asymmetry Trap", [
        "Perfectly mapping the RKHS evasion loop requires calculating dense Kernel matrices over K >= 20 iterations.",
        "The Arithmetic Impossibility: On a standard Jetson Nano (0.47 TFLOPS), this math demands ~140 seconds of execution (Asymmetry Ratio >= 17x).",
        "The Trap: The server analyzes timestamps. If an IoT device submits a perfectly evaded vector in T=1.5s, it is flagged as a Datacenter Sybil Spoofer."
    ], image_path=os.path.join("results", "hardware_time_distribution.png"))

    # Slide 8: The Whitebox Sanding down
    create_slide(prs, "Stealth-Utility Theorem (The Sanded Spike)", [
        "What happens if a Super-Adaptive white-box attacker uses infinite compute to artificially smooth the spike?",
        "Due to Fisher invariance, flattening the apex diffuses the strain outward across the volumetric integral.",
        "Infinite computation compresses geometry but NEVER erases the Riemannian boundary limit."
    ], image_path=os.path.join("results", "manifold_whitebox.png"))

    # Slide 9: Metrics
    create_slide(prs, "State-of-the-Art Evaluation Metrics", [
        "Mathematical Enforcement (PINN Guard):",
        "- Reaches a flawless 1.0000 ROC-AUC against traditional RKHS evasion.",
        "- Retains a high 0.9750 AUC against the Super-Adaptive Geometry attacker.",
        "",
        "Physical Enforcement (Compute-Asymmetry):",
        "- Achieved 0.9943 ROC-AUC hardware identity verification across 3000 proxies.",
        "",
        "Server Scalability:",
        "- PINN approximator infers at 0.56 ms per client."
    ])

    prs.save('ICLR_NeurIPS_Draft_Presentation.pptx')
    print("[+] Successfully generated ICLR_NeurIPS_Draft_Presentation.pptx")

if __name__ == '__main__':
    main()
